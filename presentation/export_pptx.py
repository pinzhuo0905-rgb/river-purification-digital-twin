"""把 16 张满版幻灯片 PNG 拼成 16:9 的 .pptx。"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

HERE = Path(__file__).parent
SLIDES_DIR = HERE / "assets" / "slides"
OUT = HERE / "项目展示.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]       # 空白版式

pngs = sorted(SLIDES_DIR.glob("slide-*.png"))
assert pngs, f"未找到幻灯片 PNG：{SLIDES_DIR}"
for png in pngs:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(png), 0, 0,
                             width=prs.slide_width, height=prs.slide_height)

prs.save(OUT)
print(f"已导出 {len(pngs)} 页 → {OUT}")
